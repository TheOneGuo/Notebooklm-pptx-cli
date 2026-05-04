#!/usr/bin/env python3
"""
NotebookLM MD → PPTX 完整流水线

一键完成：MD文件 → NotebookLM笔记本 → 横版PPT+竖版PPT同时生成 → 图片提取 → Logo遮盖

用法:
    python nb2pptx.py <md_file> [选项]

示例:
    python nb2pptx.py report.md
    python nb2pptx.py report.md --title "我的报告" --pages 40
    python nb2pptx.py report.md --output-dir ~/Desktop/output

流程:
    1. 创建 NotebookLM 笔记本
    2. 上传 MD 文件作为来源
    3. 从来源提取关键信息并重命名笔记本
    4. 生成 PPT 大纲笔记
    5. 笔记转为来源
    6. 并行生成 4 个 PPT（横版PPT1/PPT2 + 竖版PPT3/PPT4）
    7. 下载横版/竖版前后半 PPTX
    8. 生成横版 30 页和竖版 30 页 PPTX + 提取图片
    9. 遮盖 NotebookLM Logo + 绘制免责声明（仅作用于图片）
    10. 清理临时文件
"""

import argparse
import json
import os
import shutil
import socket
import subprocess
import sys
import tempfile
import time
import urllib.request
from pathlib import Path
from typing import Dict, List, Optional, Tuple

try:
    sys.stdout.reconfigure(line_buffering=True)
    sys.stderr.reconfigure(line_buffering=True)
except AttributeError:
    pass

# ═══════════════════════════════════════════════════════════════════════
# 配置
# ═══════════════════════════════════════════════════════════════════════

# NotebookLM CLI 路径（自动检测）
def find_notebooklm_cli() -> Path:
    """查找 notebooklm CLI 路径"""
    candidates = [
        Path(__file__).parent / ".venv" / "bin" / "notebooklm",
        Path.home() / ".qclaw" / "workspace-agent-ef704e53" / ".venv" / "bin" / "notebooklm",
        Path.home() / ".local" / "bin" / "notebooklm",
        Path("/usr/local/bin/notebooklm"),
    ]
    for p in candidates:
        if p.exists():
            return p
    raise FileNotFoundError("未找到 notebooklm CLI，请确保已安装")

# 默认配置
DEFAULT_LOGO_PATH = Path.home() / "Documents" / "前瞻客" / "logo和表情" / "前瞻客logo底图.png"
DEFAULT_OUTPUT_DIR = Path.home() / "Documents" / "A股研报"
DEFAULT_WAIT_CONFIG = {
    "initial_interval": 540,  # 9分钟静默等待
    "max_interval": 60,        # 后续每分钟轮询
    "timeout": 1800,           # 30分钟总超时
}
CHROME_DEBUG_PORT = int(os.environ.get("NOTEBOOKLM_CHROME_PORT", "9222"))
CHROME_PROFILE_DIR = Path(os.environ.get("NOTEBOOKLM_CHROME_PROFILE", Path.home() / ".qclaw" / "notebooklm-chrome-profile"))
COMMON_PROXY_PORTS = (7890, 7897, 7899, 7891)


def _can_connect(host: str, port: int, timeout: float = 0.4) -> bool:
    try:
        with socket.create_connection((host, port), timeout=timeout):
            return True
    except OSError:
        return False


def build_notebooklm_env() -> Dict[str, str]:
    """构建 NotebookLM CLI 环境，确保 httpx 能读到 Clash 等本地代理。"""
    env = os.environ.copy()
    has_proxy = any(env.get(k) for k in ("HTTPS_PROXY", "https_proxy", "HTTP_PROXY", "http_proxy", "ALL_PROXY", "all_proxy"))
    if not has_proxy:
        for port in COMMON_PROXY_PORTS:
            if _can_connect("127.0.0.1", port):
                proxy = f"http://127.0.0.1:{port}"
                env["HTTP_PROXY"] = proxy
                env["HTTPS_PROXY"] = proxy
                env["ALL_PROXY"] = proxy
                env.setdefault("NO_PROXY", "127.0.0.1,localhost,::1")
                print(f"🌐 检测到本地代理端口 {port}，NotebookLM CLI 将使用 {proxy}")
                break
    return env


def ensure_notebooklm_library_network_patch(cli_path: Path) -> None:
    """Patch the installed notebooklm package used by the CLI so httpx honors proxy env."""
    python_path = cli_path.parent / "python"
    if not python_path.exists():
        return

    patch_code = r'''
from pathlib import Path
import importlib.util
import re

spec = importlib.util.find_spec("notebooklm")
if not spec or not spec.submodule_search_locations:
    raise SystemExit(0)
pkg = Path(next(iter(spec.submodule_search_locations)))

def patch_file(name, replacements):
    path = pkg / name
    if not path.exists():
        return False
    text = path.read_text(encoding="utf-8")
    original = text
    for old, new in replacements:
        text = text.replace(old, new)
    if name == "_core.py":
        text = re.sub(r"DEFAULT_TIMEOUT\s*=\s*[0-9.]+", "DEFAULT_TIMEOUT = 120.0", text)
    if text != original:
        path.write_text(text, encoding="utf-8")
    return True

patch_file("_core.py", [
    ("DEFAULT_CONNECT_TIMEOUT = 10.0", "DEFAULT_CONNECT_TIMEOUT = 10.0"),
    ("httpx.AsyncClient(\n                headers=headers,\n                timeout=timeout,\n                proxy=proxy_url,\n            )",
     "httpx.AsyncClient(\n                headers=headers,\n                timeout=timeout,\n                proxy=proxy_url,\n                trust_env=True,\n            )"),
])
patch_file("auth.py", [
    ("async with httpx.AsyncClient() as client:", "async with httpx.AsyncClient(trust_env=True, timeout=120.0) as client:"),
    ("timeout=30.0,", "timeout=120.0,"),
])
patch_file("_sources.py", [
    ("httpx.AsyncClient(timeout=60.0)", "httpx.AsyncClient(timeout=120.0, trust_env=True)"),
    ("httpx.AsyncClient(timeout=300.0)", "httpx.AsyncClient(timeout=300.0, trust_env=True)"),
])
patch_file("_artifacts.py", [
    ("timeout=60.0,\n        ) as client:", "timeout=60.0,\n            trust_env=True,\n        ) as client:"),
    ("timeout=timeout,\n            ) as client:", "timeout=timeout,\n                trust_env=True,\n            ) as client:"),
])
'''
    result = subprocess.run([str(python_path), "-c", patch_code], capture_output=True, text=True, timeout=30)
    if result.returncode != 0:
        raise RuntimeError(f"NotebookLM httpx 代理补丁失败: {result.stderr.strip()}")


def chrome_debug_ready(port: int = CHROME_DEBUG_PORT) -> bool:
    try:
        with urllib.request.urlopen(f"http://127.0.0.1:{port}/json/version", timeout=2) as resp:
            return resp.status == 200
    except Exception:
        return False


def find_chrome_executable() -> Optional[Path]:
    candidates = [
        Path(os.environ["CHROME_PATH"]) if os.environ.get("CHROME_PATH") else None,
        Path("/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"),
        Path("/Applications/Google Chrome Canary.app/Contents/MacOS/Google Chrome Canary"),
        Path("/Applications/Chromium.app/Contents/MacOS/Chromium"),
    ]
    for candidate in candidates:
        if candidate and candidate.exists():
            return candidate
    return None


def ensure_chrome_debug_session(port: int = CHROME_DEBUG_PORT, profile_dir: Path = CHROME_PROFILE_DIR) -> None:
    """确保 9222 上有独立 Chrome Profile 的远程调试会话。"""
    if chrome_debug_ready(port):
        print(f"✅ Chrome Debug 端口 {port} 已可用")
        return

    chrome = find_chrome_executable()
    if not chrome:
        raise RuntimeError("未找到 Google Chrome，请先安装或通过 CHROME_PATH 指定可执行文件")

    profile_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        str(chrome),
        f"--remote-debugging-port={port}",
        f"--user-data-dir={profile_dir}",
        "--no-first-run",
        "--no-default-browser-check",
        "https://notebooklm.google.com/",
    ]
    print(f"🌐 Chrome Debug 端口 {port} 未响应，正在启动独立 Profile: {profile_dir}")
    subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

    deadline = time.time() + 20
    while time.time() < deadline:
        if chrome_debug_ready(port):
            print(f"✅ Chrome Debug 端口 {port} 已启动；如未登录，请在该 Chrome 窗口登录 Google / NotebookLM 后重跑")
            return
        time.sleep(1)

    raise RuntimeError(f"Chrome 已启动但 127.0.0.1:{port} 仍未响应，请确认远程调试窗口未被系统阻止")

# ═══════════════════════════════════════════════════════════════════════
# CLI 执行工具
# ═══════════════════════════════════════════════════════════════════════

class NotebookLM:
    """NotebookLM CLI 封装"""
    
    def __init__(self, cli_path: Path):
        self.cli = str(cli_path)
        self.env = build_notebooklm_env()
    
    def run(self, args: List[str], timeout: int = 120, check: bool = True) -> Dict:
        """执行 CLI 命令，返回 JSON 结果"""
        cmd = [self.cli] + args
        result = subprocess.run(
            cmd, 
            capture_output=True, 
            text=True, 
            timeout=timeout,
            env=self.env,
        )
        
        if check and result.returncode != 0:
            raise RuntimeError(f"命令失败: {' '.join(cmd)}\n{result.stderr}")
        
        # 尝试解析 JSON
        output = result.stdout.strip()
        if output:
            try:
                return json.loads(output)
            except json.JSONDecodeError:
                return {"output": output}
        return {}
    
    # ─── Notebook 操作 ───────────────────────────────────────────────
    
    def create_notebook(self, title: str) -> str:
        """创建笔记本，返回 notebook ID"""
        result = self.run(["notebook", "create", title, "--json"])
        return result.get("id")
    
    def rename_notebook(self, notebook_id: str, new_title: str) -> None:
        """重命名笔记本"""
        self.run(["notebook", "rename", notebook_id, new_title])
    
    def delete_notebook(self, notebook_id: str) -> None:
        """删除笔记本"""
        self.run(["notebook", "delete", notebook_id])
    
    # ─── Source 操作 ─────────────────────────────────────────────────
    
    def add_source_file(self, file_path: Path, notebook_id: str) -> Tuple[str, str]:
        """上传文件作为来源，返回 (source_id, source_title)"""
        result = self.run([
            "source", "add-file", str(file_path), 
            "-n", notebook_id, 
            "--wait", "--json"
        ])
        return result.get("id", ""), result.get("title", file_path.name)
    
    def add_source_text(self, content: str, notebook_id: str, title: str) -> str:
        """添加文本作为来源，返回 source ID"""
        # source add-text 命令格式: TITLE CONTENT [OPTIONS]
        # 注意：content 可能很长，需要小心处理
        result = self.run([
            "source", "add-text", title, content,
            "-n", notebook_id,
            "--wait", "--json"
        ], timeout=120)
        return result.get("id")
    
    def rename_source(self, source_id: str, notebook_id: str, new_title: str) -> None:
        """重命名来源"""
        self.run([
            "source", "rename", source_id, new_title,
            "-n", notebook_id
        ])
    
    def get_note_content(self, note_id: str, notebook_id: str) -> str:
        """获取笔记内容"""
        result = self.run(["notes", "get", note_id, "-n", notebook_id, "--json"], timeout=30)
        return result.get("content", "")
    
    def download_slides(
        self, 
        output_path: Path, 
        artifact_id: str, 
        notebook_id: str,
        fmt: str = "pptx"
    ) -> Path:
        """下载 PPT，返回实际文件路径"""
        result = self.run([
            "download", "slide-deck", str(output_path),
            "--artifact-id", artifact_id,
            "--format", "pptx",  # 强制使用 pptx 格式
            "-n", notebook_id,
            "--json"
        ], timeout=60)
        
        real_path = Path(result.get("output_path", output_path))
        return real_path

# ═══════════════════════════════════════════════════════════════════════
# PPT 合并工具
# ═══════════════════════════════════════════════════════════════════════

def extract_images_from_pptx(pptx_path: Path, output_dir: Path) -> List[Path]:
    """从 PPTX 提取所有图片（按页顺序）"""
    import zipfile
    import re
    import xml.etree.ElementTree as ET
    
    NS = {
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    
    images = []
    
    with zipfile.ZipFile(pptx_path, 'r') as zf:
        # 获取所有 slide 文件
        slide_files = sorted(
            [f for f in zf.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', f)],
            key=lambda x: int(re.search(r'slide(\d+)', x).group(1))
        )
        
        for slide_file in slide_files:
            slide_num = int(re.search(r'slide(\d+)', slide_file).group(1))
            slide_content = zf.read(slide_file)
            
            # 解析 .rels 文件
            rels_file = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            if rels_file not in zf.namelist():
                continue
            
            rels_content = zf.read(rels_file)
            root = ET.fromstring(rels_content)
            
            # 找到图片映射
            rid_to_target = {}
            ns_rel = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
            for rel in root.findall('.//r:Relationship', ns_rel):
                rid = rel.get('Id')
                target = rel.get('Target')
                if rid and target and 'image' in rel.get('Type', ''):
                    rid_to_target[rid] = target
            
            # 从 slide 中找主图
            slide_root = ET.fromstring(slide_content)
            main_rid = None
            max_size = 0
            
            for pic in slide_root.findall('.//p:pic', NS):
                blip = pic.find('.//a:blip', NS)
                if blip is None:
                    continue
                
                rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')
                if not rid:
                    rid = blip.get('embed')
                if not rid:
                    continue
                
                # 计算图片尺寸
                xfrm = pic.find('.//a:xfrm', NS)
                if xfrm is not None:
                    ext = xfrm.find('a:ext', NS)
                    if ext is not None:
                        try:
                            cx = int(ext.get('cx', '0'))
                            cy = int(ext.get('cy', '0'))
                            size = cx * cy
                            if size > max_size:
                                max_size = size
                                main_rid = rid
                        except (ValueError, TypeError):
                            pass
            
            if not main_rid:
                continue
            
            # 提取图片
            image_path = rid_to_target.get(main_rid)
            if not image_path:
                continue
            
            image_path = image_path.replace('../', '')
            full_path = f"ppt/{image_path}" if not image_path.startswith('ppt/') else image_path
            
            if full_path not in zf.namelist():
                continue
            
            image_data = zf.read(full_path)
            output_path = output_dir / f"P{slide_num}.png"
            
            with open(output_path, 'wb') as f:
                f.write(image_data)
            
            images.append(output_path)
    
    return images

def prepare_upload_source(md_file: Path, temp_dir: Path) -> Path:
    """复制为 ASCII 临时文件上传，避免中文路径/文件名导致 NotebookLM CLI 失败。"""
    upload_path = temp_dir / "source.md"
    shutil.copy2(md_file, upload_path)
    return upload_path


def count_pptx_slides(pptx_path: Path) -> int:
    try:
        from pptx import Presentation
    except ImportError:
        print("❌ 需要安装 python-pptx: pip install python-pptx")
        sys.exit(1)
    return len(Presentation(str(pptx_path)).slides)


def copy_images_with_expected_count(
    extracted_images: List[Path],
    output_dir: Path,
    start_page: int,
    expected_count: int,
    deck_label: str,
) -> Tuple[List[Path], List[str]]:
    """按目标页码复制图片，少页中止，多页只保留前 expected_count 页。"""
    ordered = sorted(extracted_images, key=lambda x: int(x.stem[1:]))
    if len(ordered) < expected_count:
        raise RuntimeError(f"{deck_label} 图片数 {len(ordered)} < 预期 {expected_count}，丢页中止")
    warnings = []
    if len(ordered) > expected_count:
        warnings.append(f"{deck_label} 图片数 {len(ordered)} > 预期 {expected_count}，仅保留前 {expected_count} 页")
        ordered = ordered[:expected_count]

    copied = []
    output_dir.mkdir(parents=True, exist_ok=True)
    for offset, image_path in enumerate(ordered):
        target = output_dir / f"P{start_page + offset}.png"
        shutil.copy2(image_path, target)
        copied.append(target)
    return copied, warnings


def create_pptx_from_images(images: List[Path], output_path: Path, orientation: str = "landscape") -> None:
    """用 NotebookLM 原始导出图片生成 PPTX；cover 模式保持宽高比且不留黑边。"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        from PIL import Image
    except ImportError:
        print("❌ 需要安装 python-pptx 和 Pillow")
        sys.exit(1)

    prs = Presentation()
    if orientation == "portrait":
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.333)
    else:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]
    for image_path in images:
        slide = prs.slides.add_slide(blank_layout)
        with Image.open(image_path) as img:
            img_w, img_h = img.size
        image_ratio = img_w / img_h
        slide_ratio = prs.slide_width / prs.slide_height
        if image_ratio > slide_ratio:
            pic_h = prs.slide_height
            pic_w = int(pic_h * image_ratio)
            left = int((prs.slide_width - pic_w) / 2)
            top = 0
        else:
            pic_w = prs.slide_width
            pic_h = int(pic_w / image_ratio)
            left = 0
            top = int((prs.slide_height - pic_h) / 2)
        slide.shapes.add_picture(
            str(image_path),
            left,
            top,
            width=pic_w,
            height=pic_h,
        )

    prs.save(str(output_path))
    print(f"✅ PPTX 创建完成: {output_path} ({len(images)}页, {orientation})")


def find_bad_portrait_pages(images_dir: Path, expected_pages: int, page_numbers: Optional[List[int]] = None) -> List[Dict]:
    """找出非竖版或比例明显不对的页面。"""
    try:
        from PIL import Image
    except ImportError:
        print("❌ 需要安装 Pillow: pip install Pillow")
        sys.exit(1)

    images = sorted(images_dir.glob("P*.png"), key=lambda p: int(p.stem[1:]))
    if len(images) != expected_pages:
        raise RuntimeError(f"竖版图片数 {len(images)} ≠ 预期 {expected_pages}")

    bad = []
    allowed = set(page_numbers or [])
    for image_path in images:
        page_num = int(image_path.stem[1:])
        if allowed and page_num not in allowed:
            continue
        with Image.open(image_path) as img:
            w, h = img.size
        if h <= w:
            bad.append({"page": page_num, "width": w, "height": h})
        elif h / w < 1.55:
            bad.append({"page": page_num, "width": w, "height": h})
    return bad


def validate_portrait_images(images_dir: Path, expected_pages: int, deck_label: str = "竖版", page_numbers: Optional[List[int]] = None) -> List[Dict]:
    """严格校验竖版导出图片是 9:16 纵向，不接受横版页。"""
    bad = find_bad_portrait_pages(images_dir, expected_pages, page_numbers)
    if bad:
        details = ", ".join(f"P{x['page']}({x['width']}x{x['height']})" for x in bad[:8])
        raise RuntimeError(f"竖版页面方向校验失败：{deck_label} 存在横版/非竖版页 -> {details}")
    return bad


def validate_image_count(images_dir: Path, expected_pages: int, label: str) -> List[Path]:
    images = sorted(images_dir.glob("P*.png"), key=lambda x: int(x.stem[1:]))
    if len(images) != expected_pages:
        raise RuntimeError(f"{label} 图片数 {len(images)} ≠ 预期 {expected_pages}，丢页中止")
    return images

# ═══════════════════════════════════════════════════════════════════════
# Logo 遮盖工具
# ═══════════════════════════════════════════════════════════════════════

def cover_logo_on_images(
    image_dir: Path, 
    logo_path: Path,
    logo_height: int = 24,
    margin: int = 0
) -> int:
    """
    用前瞻客 logo 遮盖 NotebookLM logo
    
    参数：
        image_dir: 图片目录
        logo_path: 前瞻客 logo 底图路径
        logo_height: logo 显示高度（px），默认 24
        margin: 距边缘边距（px），默认 0
    
    返回：
        成功处理的图片数量
    """
    try:
        from PIL import Image
    except ImportError:
        print("❌ 需要安装 Pillow: pip install Pillow")
        sys.exit(1)
    
    if not logo_path.exists():
        print(f"❌ Logo 文件不存在: {logo_path}")
        return 0
    
    # 加载并缩放 logo
    logo = Image.open(str(logo_path)).convert("RGBA")
    scale = logo_height / logo.height
    new_w = int(logo.width * scale)
    new_h = logo_height
    logo_resized = logo.resize((new_w, new_h), Image.Resampling.LANCZOS)
    
    # 处理所有图片
    images = sorted(image_dir.glob("P*.png"))
    if not images:
        print(f"⚠️ 未找到图片: {image_dir}/P*.png")
        return 0
    
    print(f"🎨 正在遮盖 {len(images)} 张图片的 NotebookLM logo...")
    success = 0
    
    for img_path in images:
        try:
            img = Image.open(str(img_path)).convert("RGBA")
            img_w, img_h = img.size
            
            # 右下角对齐
            pos_x = img_w - new_w - margin
            pos_y = img_h - new_h - margin
            
            # 粘贴 logo
            img.paste(logo_resized, (pos_x, pos_y), logo_resized)
            
            # 保存（转为 RGB 去掉 alpha 通道）
            img.convert("RGB").save(str(img_path), optimize=True)
            success += 1
            
        except Exception as e:
            print(f"   ✗ {img_path.name}: {e}")
    
    print(f"✅ Logo 遮盖完成: {success}/{len(images)}")
    return success

# ═══════════════════════════════════════════════════════════════════════
# 免责声明绘制工具
# ═══════════════════════════════════════════════════════════════════════

def draw_disclaimer(
    image_dir: Path,
    orientation: str = "landscape",
    font_size: int = 11,
    text_color: tuple = (170, 170, 170),
    bg_color: tuple = (30, 30, 30, 180),
    margin: int = 8
) -> int:
    """
    在图片底部绘制标准股市免责声明

    参数：
        image_dir: 图片目录
        orientation: "landscape" (横版) 或 "portrait" (竖版)
        font_size: 字号，默认 11
        text_color: 文字颜色 RGB，默认 #AAAAAA 暗灰
        bg_color: 背景颜色 RGBA，默认半透明深色
        margin: 距底部边距（px），默认 8

    返回：
        成功处理的图片数量
    """
    try:
        from PIL import Image, ImageDraw, ImageFont
    except ImportError:
        print("❌ 需要安装 Pillow: pip install Pillow")
        return 0

    disclaimer_lines = [
        "市场有风险，决策需独立；",
        "股市有风险，入市需谨慎。"
    ]

    # 查找中文字体
    font = None
    font_paths = [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Medium.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    for fp in font_paths:
        if Path(fp).exists():
            try:
                font = ImageFont.truetype(fp, font_size)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()
        print("   ⚠️ 未找到中文字体，使用默认字体（免责声明可能显示异常）")

    images = sorted(image_dir.glob("P*.png"))
    if not images:
        print(f"⚠️ 未找到图片: {image_dir}/P*.png")
        return 0

    print(f"📝 正在绘制免责声明（{len(images)} 张图片）...")
    success = 0

    for img_path in images:
        try:
            img = Image.open(str(img_path)).convert("RGBA")
            img_w, img_h = img.size
            draw = ImageDraw.Draw(img)

            # 计算文字区域
            line_heights = []
            line_widths = []
            for line in disclaimer_lines:
                bbox = draw.textbbox((0, 0), line, font=font)
                line_widths.append(bbox[2] - bbox[0])
                line_heights.append(bbox[3] - bbox[1])

            text_w = max(line_widths)
            text_h = sum(line_heights) + 4  # 行间距
            line_spacing = 4

            # 横版和竖版都只放在页面正中底部，避免右下角重复免责声明。
            text_x = (img_w - text_w) // 2
            text_y = img_h - text_h - margin - 4

            # 绘制半透明背景
            pad = 6
            bg_x1 = text_x - pad
            bg_y1 = text_y - pad
            bg_x2 = text_x + text_w + pad
            bg_y2 = text_y + text_h + pad

            # 创建背景遮罩
            bg_overlay = Image.new("RGBA", img.size, (0, 0, 0, 0))
            bg_draw = ImageDraw.Draw(bg_overlay)
            bg_draw.rectangle([bg_x1, bg_y1, bg_x2, bg_y2], fill=bg_color)
            img = Image.alpha_composite(img, bg_overlay)
            draw = ImageDraw.Draw(img)

            # 绘制文字
            y_offset = text_y
            for i, line in enumerate(disclaimer_lines):
                line_w = line_widths[i]
                lx = text_x + (text_w - line_w) // 2
                draw.text((lx, y_offset), line, fill=text_color, font=font)
                y_offset += line_heights[i] + line_spacing

            img.convert("RGB").save(str(img_path), optimize=True)
            success += 1

        except Exception as e:
            print(f"   ✗ {img_path.name}: {e}")

    print(f"✅ 免责声明绘制完成: {success}/{len(images)}")
    return success


# ═══════════════════════════════════════════════════════════════════════
# 主流程
# ═══════════════════════════════════════════════════════════════════════

def check_dependencies() -> List[str]:
    """检查运行依赖，返回缺失项列表"""
    missing = []
    try:
        import pptx  # noqa: F401
    except ImportError:
        missing.append("python-pptx (pip install python-pptx)")
    try:
        from PIL import Image  # noqa: F401
    except ImportError:
        missing.append("Pillow (pip install Pillow)")
    return missing


def main(
    md_file: Path,
    title: Optional[str] = None,
    pages: int = 30,
    output_dir: Optional[Path] = None,
    logo_path: Optional[Path] = None,
    keep_temp: bool = False,
    wait_config: Optional[Dict] = None,
    check_chrome: bool = True,
):
    """
    完整流水线：MD → NotebookLM → PPTX → 图片 → Logo 遮盖
    """
    # ─── 前置检查 ────────────────────────────────────────────────────
    missing = check_dependencies()
    if missing:
        print("❌ 缺少运行依赖:")
        for m in missing:
            print(f"   - {m}")
        sys.exit(1)

    if md_file.stat().st_size == 0:
        print(f"❌ 输入文件为空: {md_file}")
        sys.exit(1)
    if pages <= 0 or pages % 2 != 0:
        print("❌ --pages 必须是正偶数；当前流程需要前半/后半各一份 PPT")
        sys.exit(1)
    if pages != 30:
        print(f"⚠️ 当前视频拆分流程按 30 页验收；本次使用 {pages} 页，请确认这是预期行为")

    # 参数处理
    title = title or md_file.stem
    logo_path = logo_path or DEFAULT_LOGO_PATH
    wait_config = wait_config or DEFAULT_WAIT_CONFIG

    if not logo_path.exists():
        print(f"⚠️ Logo 文件不存在: {logo_path}，将跳过 Logo 遮盖")
        logo_path = None

    print("="*60)
    print(f"📄 输入文件: {md_file} ({md_file.stat().st_size} bytes)")
    print(f"📝 初始标题: {title}")
    print(f"📊 目标页数: {pages}")
    print(f"🎨 Logo: {logo_path or '(跳过)'}")
    print(f"📁 输出目录: {output_dir or '默认'}")
    print("="*60)

    # 初始化
    cli_path = find_notebooklm_cli()
    ensure_notebooklm_library_network_patch(cli_path)
    if check_chrome:
        ensure_chrome_debug_session()
    nb = NotebookLM(cli_path)
    
    notebook_id = None
    temp_dir = Path(tempfile.mkdtemp(prefix="nb2pptx_"))
    print(f"📂 临时目录: {temp_dir}")
    notebook_title = title  # 初始化为 title，后续会被步骤3重命名
    run_id = time.strftime("%Y%m%d_%H%M%S")
    run_start = time.time()
    step_times = []  # [(step_name, elapsed_sec)]

    def step_timer(step_name: str):
        """记录上一步耗时并开始新步骤"""
        elapsed = time.time() - run_start
        if step_times:
            prev_name, prev_start = step_times[-1]
            step_elapsed = elapsed - prev_start
            step_times[-1] = (prev_name, step_elapsed)
        step_times.append((step_name, elapsed))
        return elapsed
    
    try:
        # ─── Step 1: 创建笔记本 ───────────────────────────────────────
        step_timer("1.创建笔记本")
        print("\n[1/10] 创建 NotebookLM 笔记本...")
        notebook_id = nb.create_notebook(title)
        if not notebook_id:
            raise RuntimeError("创建笔记本失败：返回空 ID")
        print(f"   ✅ 笔记本 ID: {notebook_id}")

        # ─── Step 2: 上传 MD 作为来源并重命名为 a ────────────────────
        step_timer("2.上传来源")
        print("\n[2/10] 上传 MD 文件作为来源...")
        upload_md_file = prepare_upload_source(md_file, temp_dir)
        print(f"   → 使用 ASCII 临时文件上传，规避中文路径问题: {upload_md_file.name}")
        source_a_id, source_a_title = nb.add_source_file(upload_md_file, notebook_id)
        if not source_a_id:
            raise RuntimeError("上传来源失败：返回空 source ID")
        print(f"   ✅ 来源上传成功: {source_a_title} ({source_a_id})")
        
        # 重命名为 a
        print("   → 重命名为 'a'...")
        nb.rename_source(source_a_id, notebook_id, "a")
        print(f"   ✅ 来源 A: a ({source_a_id})")
        
        # ─── Step 3: 从来源提取关键信息并重命名笔记本 ──────────────────
        step_timer("3.提取主题")
        print("\n[3/10] 从来源提取关键信息并重命名笔记本...")
        extract_prompt = '请从来源 a 中提取报告的核心主题，用一句话概括（不超过20字），仅返回主题名称，不要任何其他内容。不要添加"answer:"前缀。'
        result = nb.run([
            "ask", extract_prompt,
            "-n", notebook_id,
            "--source", source_a_id
        ], timeout=60)
        
        # 从结果中提取标题
        output_text = result.get("output", "").strip()
        if not output_text:
            output_text = str(result).strip()
        
        # 清理输出：去除 "answer:" 或 "answer：" 前缀，以及多余内容
        output_text = output_text.replace("answer:", "").replace("answer：", "").strip()
        
        notebook_title = output_text.split('\n')[0][:20] if output_text else title
        
        # 重命名笔记本
        nb.rename_notebook(notebook_id, notebook_title)
        print(f"   ✅ 笔记本重命名为: {notebook_title}")
        
        # ─── Step 4: 自动生成 PPT 大纲（带风格要求）────────────────────
        step_timer("4.生成大纲")
        print("\n[4/10] 自动生成 PPT 大纲（带风格要求）...")
        
        # 构建提示词（包含详细的风格要求，强制中文）
        prompt = f'''你现在是一位资深的商业咨询顾问，根据来源a的研报内容，精确提炼核心内容生成一份{pages}页的PPT大纲。

【语言要求】所有内容必须使用中文，包括标题、要点、章节名称等。严禁出现英文。

【内容要求】完全按照研报的内容进行设计，每页都要求突出核心重点，主要内容不缺失。研报中提到的所有公司名称、股票代码、核心数据必须在大纲中明确体现，不得遗漏。每页包含：1) 明确的页面标题 2) 3-5个核心要点 3) 逻辑连贯，适合演示。

【风格要求】PPT设计风格采用高端商务黑金/暗黑机械风格。背景使用深邃黑与暗灰色渐变，营造沉浸式暗室感。主体图表必须呈现 3D 拟物化的哑光红铜/古铜金属材质。图表边缘需带有亮金色流光特效。整体画面需传达出极度沉稳、权威、严谨与精密的质感。

【页码要求】第1页必须是标题页；第16页必须承接第15页，从来源b大纲的第16页内容继续，不能重新开始。

【免责声明】不要在大纲或 PPT 画面里自行添加任何免责声明、页脚免责声明或右下角免责声明；程序会在导出的图片正中底部统一绘制。'''
        
        note_title = f"{pages}页PPT大纲"
        
        # 使用 ask --save-as-note 生成大纲并保存
        # 注意：生成大纲可能需要 3-5 分钟，设置 300s 超时
        result = nb.run([
            "ask", prompt,
            "-n", notebook_id,
            "--source", source_a_id,
            "--save-as-note",
            "--note-title", note_title
        ], timeout=900)  # 大报告可能超过 58KB，允许更长时间
        
        print(f"   ✅ 大纲笔记已生成")
        
        # 查找笔记 ID
        notes_result = nb.run(["notes", "list", "-n", notebook_id, "--json"], timeout=30)
        note_id = None
        for note in notes_result.get("notes", []):
            if note.get("title") == note_title:
                note_id = note.get("id")
                break
        
        if not note_id:
            raise RuntimeError(f"未找到笔记：{note_title}")
        
        print(f"   ✅ 笔记 ID: {note_id}")
        
        # ─── Step 5: 笔记转为来源 ───────────────────────────────────
        step_timer("5.笔记转来源")
        print("\n[5/10] 将大纲笔记转为来源 B...")
        # 获取笔记内容
        note_content = nb.get_note_content(note_id, notebook_id)
        if not note_content or len(note_content.strip()) < 50:
            raise RuntimeError(f"笔记内容为空或过短（{len(note_content)} 字符），无法生成 PPT")
        print(f"   📝 笔记内容: {len(note_content)} 字符")

        # 添加为来源
        source_b_id = nb.add_source_text(note_content, notebook_id, note_title)
        if not source_b_id:
            raise RuntimeError("添加来源 B 失败：返回空 source ID")
        # 重命名为 b
        nb.rename_source(source_b_id, notebook_id, "b")
        print(f"   ✅ 来源 B: {source_b_id}")
        
        # ─── Step 6: 并行生成 PPT（横版 PPT1/PPT2 + 竖版 PPT3/PPT4）──────────────────────
        step_timer("6.并行生成PPT")
        print("\n[6/10] 并行生成 PPT（横版+竖版）...")
        pages_per_deck = pages // 2
        
        # 获取生成前的 artifacts 列表（用于排除已存在的）
        artifacts_result = nb.run(["artifact", "list", "-n", notebook_id, "--json"], timeout=30)
        initial_artifacts = artifacts_result.get("artifacts", [])
        initial_artifact_ids = {a.get("id") for a in initial_artifacts}
        print(f"   → 已有 {len(initial_artifacts)} 个 artifacts")
        
        # ═══════════════════════════════════════════════════════════════════
        # PPT 提示词（公共模板 + 4 个任务）
        # ═══════════════════════════════════════════════════════════════════

        def build_ppt_prompt(orientation: str, page_start: int, page_end: int, is_first_half: bool) -> str:
            """构建 PPT 生成提示词

            Args:
                orientation: "landscape" (横版16:9) 或 "portrait" (竖版9:16)
                page_start: 起始页码
                page_end: 结束页码
                is_first_half: 是否为前半段（影响"先制作"vs"制作"措辞）
            """
            if orientation == "landscape":
                layout_rule = "横版（16:9）。页面宽度大于高度，适合PC端展示和宽屏投影。"
                label = "横版"
            else:
                layout_rule = (
                    "必须原生设计为竖版（9:16）纵向布局，页面高度明显大于宽度（高宽比约为1.77:1），"
                    "适合手机端短视频展示。必须重新按竖版排版和设计，不得把横版页面、横版图片、横向四栏仪表盘、"
                    "横向宽表或横向长图压缩进竖版画布；也不得把横版内容等比缩小后上下留黑边/暗色填充。"
                    "所有内容必须采用移动端原生纵向叙事：大标题在上方，核心结论分层向下展开，图表改为单列、双层、"
                    "上下堆叠卡片、纵向流程或竖向对比。宁可减少同屏元素，也不能压缩文字、人物、图表或图片。"
                )
                label = "竖版"

            content_verb = "先制作" if is_first_half else "制作"
            style_ref = "来源b的大纲" if is_first_half else "来源b里的大纲"
            consistency = "" if is_first_half else f"，要求和前{page_start - 1}页保持完全一致的风格和语言"

            return f'''[LANGUAGE: CHINESE ONLY - 禁止使用任何英文]

根据来源a的分析内容，按照{style_ref}制作PPT。这份PPT总共包含{pages}页。

【强制语言要求】
- 所有页面必须100%使用中文
- 标题、正文、图表标签、页脚、按钮文字必须全部是中文
- 严禁出现任何英文单词或字母
- 如果必须使用专有名词，请用中文翻译

【版式要求】{layout_rule}

【风格要求】采用高端商务黑金/暗黑机械风格。背景使用深邃黑与暗灰色渐变，营造沉浸式暗室感。主体图表（如结构图、流程图）必须呈现 3D 拟物化的哑光红铜/古铜金属材质，避免亮金色。图表边缘需带有亮金色流光特效。整体画面需传达出极度沉稳、权威、严谨与精密的质感。

【页码闭环】PPT 1 和 PPT 3 都必须覆盖来源b大纲的第1-{pages_per_deck}页，其中第1页是标题页；PPT 2 和 PPT 4 都必须覆盖第{pages_per_deck + 1}-{pages}页，其中第{pages_per_deck + 1}页是第{pages_per_deck}页之后的续篇，不能重新做标题页。横版和竖版每一页的主题、标题含义、公司/股票/数据必须一一对应，只允许版式和视觉设计不同。

【免责声明】不要自行添加任何页脚、免责声明或右下角小字；程序会在导出的图片正中底部统一绘制。

【内容要求】要求插图丰富，确保每个中文字不要出错，字体清晰。研报中提到的所有公司名称、股票代码、核心数据必须在PPT中明确体现，不得遗漏。{content_verb}来源b要求的第{page_start}-{page_end}页{consistency}。

【PPT标识】这是{label}PPT，包含第{page_start}-{page_end}页。'''

        # 4 个 PPT 任务
        ppt1_prompt = build_ppt_prompt("landscape", 1, pages_per_deck, is_first_half=True)
        ppt2_prompt = build_ppt_prompt("landscape", pages_per_deck + 1, pages, is_first_half=False)
        ppt3_prompt = build_ppt_prompt("portrait", 1, pages_per_deck, is_first_half=True)
        ppt4_prompt = build_ppt_prompt("portrait", pages_per_deck + 1, pages, is_first_half=False)

        # ═══════════════════════════════════════════════════════════════════
        # 并行提交 4 个生成任务
        # ═══════════════════════════════════════════════════════════════════
        
        def submit_and_track(task_name, prompt, excluded_ids):
            """提交生成任务并追踪 artifact ID"""
            print(f"\n   → 提交 {task_name} 生成任务...")
            result = nb.run([
                "generate", "slide-deck",
                "--source", source_b_id,
                "--instructions", prompt,
                "--format", "detailed_deck",
                "--no-wait",
                "-n", notebook_id,
                "--json"
            ], timeout=30)
            
            artifact_id = result.get("artifact_id") or result.get("task_id")
            for retry in range(10):
                time.sleep(5 + retry * 3)
                artifacts_resp = nb.run(["artifact", "list", "-n", notebook_id, "--json"], timeout=30)
                artifacts_list = artifacts_resp.get("artifacts", [])
                artifact_ids = {a.get("id") for a in artifacts_list}

                if artifact_id and artifact_id in artifact_ids:
                    status = next(a.get("status") for a in artifacts_list if a.get("id") == artifact_id)
                    print(f"   ✅ {task_name} artifact 已绑定: {artifact_id[:8]}... (状态: {status})")
                    return artifact_id, excluded_ids | {artifact_id}
                
                new_artifacts = [
                    a for a in artifacts_list
                    if a.get("id") not in excluded_ids
                ]
                
                if len(new_artifacts) >= 1:
                    artifact_id = new_artifacts[0].get("id")
                    status = new_artifacts[0].get("status")
                    print(f"   ✅ {task_name} artifact 已创建: {artifact_id[:8]}... (状态: {status})")
                    # 返回更新后的 excluded_ids（包含新发现的 artifact）
                    new_excluded = excluded_ids | {a.get("id") for a in artifacts_list}
                    return artifact_id, new_excluded
                else:
                    print(f"   ⚠️ 第{retry+1}次查询未找到 {task_name} artifact，{'重试中...' if retry < 9 else '将在下载后检查内容'}")
            
            if not artifact_id:
                print(f"   ⚠️ {task_name} artifact ID 获取失败（10次重试均未找到），将在下载阶段尝试方案C")
                return None, excluded_ids
            print(f"   ⚠️ {task_name} 暂未出现在 artifact 列表，保留返回 ID 继续等待: {artifact_id}")
            return artifact_id, excluded_ids

        def submit_single_portrait_page(page_number: int, excluded_ids: set) -> str:
            prompt = f'''[LANGUAGE: CHINESE ONLY - 禁止使用任何英文]

请只重新生成 1 页竖版 PPT：对应总大纲的第 {page_number} 页。

【强制要求】
- 只生成 1 页，不要生成标题页、目录页或额外页面。
- 必须是原生 9:16 竖版设计，页面高度明显大于宽度。
- 不得把横版页面、横版图片、横向宽图、横向仪表盘压缩进竖版画布。
- 不得使用上下黑边、暗色填充或等比缩小横版图的方式冒充竖版。
- 内容必须与来源b第 {page_number} 页的主题、标题含义、公司/股票/数据一致。
- 采用移动端竖向叙事：标题在上、核心结论纵向分层、图表单列或上下堆叠。
- 不要自行添加任何页脚或免责声明，程序会统一绘制。

【PPT标识】这是竖版单页重做，页码 P{page_number}。'''
            artifact_id, _ = submit_and_track(f"P{page_number} 竖版单页重做", prompt, excluded_ids)
            if not artifact_id:
                raise RuntimeError(f"P{page_number} 竖版单页重做未能绑定 artifact")
            return artifact_id

        def wait_for_artifact_ready(label: str, artifact_id: str, timeout_seconds: float) -> None:
            start = time.time()
            while True:
                if time.time() - start > timeout_seconds:
                    raise RuntimeError(f"{label} 等待完成超时（{timeout_seconds}s）")
                result = nb.run(["artifact", "list", "-n", notebook_id, "--json"], timeout=30)
                artifact_map = {a.get("id"): a for a in result.get("artifacts", [])}
                artifact = artifact_map.get(artifact_id)
                if artifact:
                    status_code = artifact.get("status", 0)
                    if status_code == 3:
                        print(f"   ✅ {label} 已完成")
                        return
                    if status_code == 4:
                        raise RuntimeError(f"{label} 生成失败")
                    print(f"   → {label} 当前状态: {STATUS_MAP.get(status_code, status_code)}")
                time.sleep(60)
        
        # 按顺序提交 4 个任务，每个都追踪 artifact ID
        # PPT 1（横版前半）
        ppt1_artifact_id, excluded_after_ppt1 = submit_and_track(
            "PPT 1（横版前半）", ppt1_prompt, initial_artifact_ids
        )
        
        # PPT 2（横版后半）
        ppt2_artifact_id, excluded_after_ppt2 = submit_and_track(
            "PPT 2（横版后半）", ppt2_prompt, excluded_after_ppt1
        )
        
        # PPT 3（竖版前半）
        ppt3_artifact_id, excluded_after_ppt3 = submit_and_track(
            "PPT 3（竖版前半）", ppt3_prompt, excluded_after_ppt2
        )
        
        # PPT 4（竖版后半）
        ppt4_artifact_id, excluded_after_ppt4 = submit_and_track(
            "PPT 4（竖版后半）", ppt4_prompt, excluded_after_ppt3
        )
        
        # ─── 等待 4 个 PPT 完成 ──────────────────────────────────────
        print("\n   ⏳ 等待 PPT 生成完成...")
        
        # 状态码映射
        STATUS_MAP = {0: "pending", 1: "processing", 2: "ready", 3: "completed", 4: "failed"}
        
        start_time = time.time()
        timeout = wait_config.get("timeout", 1800)  # 4个PPT，超时改为30分钟
        initial_interval = wait_config.get("initial_interval", 540)
        max_interval = wait_config.get("max_interval", 60)
        
        # 初始静默等待
        if initial_interval > 0:
            print(f"   → 初始静默等待 {initial_interval}s...")
            time.sleep(initial_interval)
        
        poll_count = 0
        ppt1_completed = False
        ppt2_completed = False
        ppt3_completed = False
        ppt4_completed = False
        
        while not (ppt1_completed and ppt2_completed and ppt3_completed and ppt4_completed):
            elapsed = time.time() - start_time
            if elapsed > timeout:
                raise RuntimeError(f"PPT 生成超时（{timeout}s）")
            
            # 查询 artifact 列表
            result = nb.run(["artifact", "list", "-n", notebook_id, "--json"], timeout=30)
            artifacts = result.get("artifacts", [])
            
            # 创建 ID -> artifact 的映射
            artifact_map = {a.get("id"): a for a in artifacts}

            def bind_missing(current_id, task_label):
                if current_id:
                    return current_id
                assigned_ids = {x for x in [ppt1_artifact_id, ppt2_artifact_id, ppt3_artifact_id, ppt4_artifact_id] if x}
                candidates = [
                    a for a in artifacts
                    if a.get("kind") == "slide_deck"
                    and a.get("id") not in initial_artifact_ids
                    and a.get("id") not in assigned_ids
                ]
                if not candidates:
                    return None
                picked = candidates[0].get("id")
                print(f"   ✅ {task_label} 迟到 artifact 已绑定: {picked[:8]}...")
                return picked

            ppt1_artifact_id = bind_missing(ppt1_artifact_id, "PPT 1（横版前半）")
            ppt2_artifact_id = bind_missing(ppt2_artifact_id, "PPT 2（横版后半）")
            ppt3_artifact_id = bind_missing(ppt3_artifact_id, "PPT 3（竖版前半）")
            ppt4_artifact_id = bind_missing(ppt4_artifact_id, "PPT 4（竖版后半）")
            
            # 检查 PPT 1 状态
            if not ppt1_completed and ppt1_artifact_id and ppt1_artifact_id in artifact_map:
                status_code = artifact_map[ppt1_artifact_id].get("status", 0)
                if status_code == 3:
                    print(f"   ✅ PPT 1（横版前半）已完成")
                    ppt1_completed = True
                elif status_code == 4:
                    raise RuntimeError("PPT 1（横版前半）生成失败")
            
            # 检查 PPT 2 状态
            if not ppt2_completed and ppt2_artifact_id and ppt2_artifact_id in artifact_map:
                status_code = artifact_map[ppt2_artifact_id].get("status", 0)
                if status_code == 3:
                    print(f"   ✅ PPT 2（横版后半）已完成")
                    ppt2_completed = True
                elif status_code == 4:
                    raise RuntimeError("PPT 2（横版后半）生成失败")
            
            # 检查 PPT 3 状态
            if not ppt3_completed and ppt3_artifact_id and ppt3_artifact_id in artifact_map:
                status_code = artifact_map[ppt3_artifact_id].get("status", 0)
                if status_code == 3:
                    print(f"   ✅ PPT 3（竖版前半）已完成")
                    ppt3_completed = True
                elif status_code == 4:
                    raise RuntimeError("PPT 3（竖版前半）生成失败")
            
            # 检查 PPT 4 状态
            if not ppt4_completed and ppt4_artifact_id and ppt4_artifact_id in artifact_map:
                status_code = artifact_map[ppt4_artifact_id].get("status", 0)
                if status_code == 3:
                    print(f"   ✅ PPT 4（竖版后半）已完成")
                    ppt4_completed = True
                elif status_code == 4:
                    raise RuntimeError("PPT 4（竖版后半）生成失败")
            
            # 如果全部完成，退出循环
            if ppt1_completed and ppt2_completed and ppt3_completed and ppt4_completed:
                break
            
            poll_count += 1
            s1 = STATUS_MAP.get(artifact_map.get(ppt1_artifact_id, {}).get("status", 0), "unknown") if ppt1_artifact_id else "N/A"
            s2 = STATUS_MAP.get(artifact_map.get(ppt2_artifact_id, {}).get("status", 0), "unknown") if ppt2_artifact_id else "N/A"
            s3 = STATUS_MAP.get(artifact_map.get(ppt3_artifact_id, {}).get("status", 0), "unknown") if ppt3_artifact_id else "N/A"
            s4 = STATUS_MAP.get(artifact_map.get(ppt4_artifact_id, {}).get("status", 0), "unknown") if ppt4_artifact_id else "N/A"
            print(f"   → 第{poll_count}次轮询：横1={s1}, 横2={s2}, 竖1={s3}, 竖2={s4}")
            
            time.sleep(max_interval)
        
        # 方案C：如果无法记录ID，下载后检查内容
        if not all([ppt1_artifact_id, ppt2_artifact_id, ppt3_artifact_id, ppt4_artifact_id]):
            print("\n   ⚠️ 无法通过立即查询记录部分 artifact ID，将下载后检查内容（方案C）")
        
        print(f"\n   ✅ PPT 1（横版前半）: {ppt1_artifact_id}")
        print(f"   ✅ PPT 2（横版后半）: {ppt2_artifact_id}")
        print(f"   ✅ PPT 3（竖版前半）: {ppt3_artifact_id}")
        print(f"   ✅ PPT 4（竖版后半）: {ppt4_artifact_id}")
        
        # ─── Step 7: 下载 PPT ─────────────────────────────────────
        step_timer("7.下载PPT")
        print("\n[7/10] 下载 PPT...")
        
        def download_with_fallback(task_label, artifact_id, temp_name, excluded_ids_set):
            """下载 PPT，如果 ID 缺失则通过内容识别"""
            if artifact_id:
                path = nb.download_slides(temp_dir / temp_name, artifact_id, notebook_id)
                print(f"   ✅ {task_label}: {path.name}")
                excluded_ids_set.add(artifact_id)
                return path
            
            # 方案C：通过内容识别
            print(f"   ⚠️ {task_label} ID 缺失，执行方案C...")
            result = nb.run(["artifact", "list", "-n", notebook_id, "--json"], timeout=30)
            artifacts = result.get("artifacts", [])
            
            candidates = [
                a for a in artifacts
                if a.get("status") == 3
                and a.get("kind") == "slide_deck"
                and a.get("id") not in excluded_ids_set
            ]
            
            if not candidates:
                raise RuntimeError(f"未找到 {task_label} 的 artifact")
            
            picked = candidates[0]
            picked_id = picked.get("id")
            path = nb.download_slides(temp_dir / temp_name, picked_id, notebook_id)
            print(f"   ✅ {task_label} (方案C): {path.name}")
            excluded_ids_set.add(picked_id)
            return path
        
        excluded_ids = set(initial_artifact_ids)
        
        # 下载横版 PPT1 + PPT2
        print("   → 下载横版 PPT...")
        pptx_h1 = download_with_fallback("PPT 1（横版前半）", ppt1_artifact_id, "ppt_h1.pptx", excluded_ids)
        pptx_h2 = download_with_fallback("PPT 2（横版后半）", ppt2_artifact_id, "ppt_h2.pptx", excluded_ids)
        
        # 下载竖版 PPT3 + PPT4
        print("   → 下载竖版 PPT...")
        pptx_v1 = download_with_fallback("PPT 3（竖版前半）", ppt3_artifact_id, "ppt_v1.pptx", excluded_ids)
        pptx_v2 = download_with_fallback("PPT 4（竖版后半）", ppt4_artifact_id, "ppt_v2.pptx", excluded_ids)
        
        # ─── Step 8: 生成 30 页 PPTX + 提取图片 ──────────────────────
        step_timer("8.合并+图片")
        print("\n[8/10] 生成 30 页横版/竖版 PPTX，并提取视频图片...")
        
        # 创建输出目录（尊重 --output-dir 参数）
        final_output_dir = output_dir if output_dir else DEFAULT_OUTPUT_DIR / notebook_title
        final_output_dir.mkdir(parents=True, exist_ok=True)
        print(f"   📁 输出目录: {final_output_dir}")
        
        final_pptx_h = final_output_dir / f"{notebook_title}_横版.pptx"
        final_pptx_v = final_output_dir / f"{notebook_title}_竖版.pptx"

        warnings = []

        print("\n   📐 处理横版 PPT（PPT1 P1-P15 + PPT2 P16-P30）...")
        raw_images_h = temp_dir / "raw_images_landscape"
        raw_images_h.mkdir(exist_ok=True)
        h1_extract_dir = temp_dir / "extract_h1"
        h2_extract_dir = temp_dir / "extract_h2"
        h1_extract_dir.mkdir(exist_ok=True)
        h2_extract_dir.mkdir(exist_ok=True)
        images_h1 = extract_images_from_pptx(pptx_h1, h1_extract_dir)
        _, warnings_h1 = copy_images_with_expected_count(images_h1, raw_images_h, 1, pages_per_deck, "横版前半")
        images_h2 = extract_images_from_pptx(pptx_h2, h2_extract_dir)
        _, warnings_h2 = copy_images_with_expected_count(images_h2, raw_images_h, pages_per_deck + 1, pages - pages_per_deck, "横版后半")
        warnings.extend(warnings_h1 + warnings_h2)
        raw_all_images_h = validate_image_count(raw_images_h, pages, "横版")
        create_pptx_from_images(raw_all_images_h, final_pptx_h, orientation="landscape")
        if count_pptx_slides(final_pptx_h) != pages:
            raise RuntimeError(f"横版 PPTX 页数 {count_pptx_slides(final_pptx_h)} ≠ 预期 {pages}")

        print("\n   📱 处理竖版 PPT（PPT3 P1-P15 + PPT4 P16-P30）...")
        raw_images_v = temp_dir / "raw_images_portrait"
        raw_images_v.mkdir(exist_ok=True)
        v1_extract_dir = temp_dir / "extract_v1"
        v2_extract_dir = temp_dir / "extract_v2"
        v1_extract_dir.mkdir(exist_ok=True)
        v2_extract_dir.mkdir(exist_ok=True)
        images_v1 = extract_images_from_pptx(pptx_v1, v1_extract_dir)
        _, warnings_v1 = copy_images_with_expected_count(images_v1, raw_images_v, 1, pages_per_deck, "竖版前半")
        images_v2 = extract_images_from_pptx(pptx_v2, v2_extract_dir)
        _, warnings_v2 = copy_images_with_expected_count(images_v2, raw_images_v, pages_per_deck + 1, pages - pages_per_deck, "竖版后半")
        warnings.extend(warnings_v1 + warnings_v2)
        raw_all_images_v = validate_image_count(raw_images_v, pages, "竖版")
        bad_portrait_pages = find_bad_portrait_pages(raw_images_v, pages)
        if bad_portrait_pages:
            rerender_excluded_ids = set(initial_artifact_ids) | {
                x for x in [ppt1_artifact_id, ppt2_artifact_id, ppt3_artifact_id, ppt4_artifact_id] if x
            }
            print(f"   ⚠️ 竖版发现 {len(bad_portrait_pages)} 页横版/非竖版，将逐页单独重做")
            for item in bad_portrait_pages:
                page_number = item["page"]
                print(f"      ↻ 重做竖版 P{page_number}（原尺寸 {item['width']}x{item['height']}）")
                rerender_artifact_id = submit_single_portrait_page(page_number, rerender_excluded_ids)
                rerender_excluded_ids.add(rerender_artifact_id)
                wait_for_artifact_ready(f"P{page_number} 竖版单页重做", rerender_artifact_id, timeout)
                rerender_pptx = nb.download_slides(temp_dir / f"portrait_rerender_P{page_number}.pptx", rerender_artifact_id, notebook_id)
                rerender_extract_dir = temp_dir / f"extract_portrait_rerender_P{page_number}"
                rerender_extract_dir.mkdir(exist_ok=True)
                rerender_images = extract_images_from_pptx(rerender_pptx, rerender_extract_dir)
                copied, rerender_warnings = copy_images_with_expected_count(
                    rerender_images,
                    raw_images_v,
                    page_number,
                    1,
                    f"竖版 P{page_number} 单页重做",
                )
                warnings.extend(rerender_warnings)
                replacement_bad = find_bad_portrait_pages(raw_images_v, pages, page_numbers=[page_number])
                if replacement_bad:
                    detail = replacement_bad[0]
                    raise RuntimeError(
                        f"竖版 P{page_number} 单页重做后仍非竖版: {detail['width']}x{detail['height']}"
                    )
                print(f"      ✅ P{page_number} 已替换为 {copied[0].name}")
            raw_all_images_v = validate_image_count(raw_images_v, pages, "竖版")
        validate_portrait_images(raw_images_v, pages, "竖版")
        create_pptx_from_images(raw_all_images_v, final_pptx_v, orientation="portrait")
        if count_pptx_slides(final_pptx_v) != pages:
            raise RuntimeError(f"竖版 PPTX 页数 {count_pptx_slides(final_pptx_v)} ≠ 预期 {pages}")

        # ── 横版视频图片后处理 ───────────────────────────────────────
        print("\n   📐 生成横版视频图片（16:9）...")
        images_dir_h = final_output_dir / "images_landscape"
        if images_dir_h.exists():
            shutil.rmtree(images_dir_h)
        images_dir_h.mkdir(exist_ok=True)
        for image_path in raw_all_images_h:
            shutil.copy2(image_path, images_dir_h / image_path.name)
        all_images_h = validate_image_count(images_dir_h, pages, "横版")

        if logo_path:
            print("      → 遮盖横版 NotebookLM Logo...")
            cover_h = cover_logo_on_images(images_dir_h, logo_path)
            if cover_h != len(all_images_h):
                print(f"   ⚠️ 横版部分图片遮盖失败: {cover_h}/{len(all_images_h)}")
        else:
            print("      ⏭️ 跳过 Logo 遮盖（Logo 文件不存在）")
        print("      → 绘制横版免责声明...")
        draw_disclaimer(images_dir_h, orientation="landscape")

        # ── 竖版视频图片后处理 ───────────────────────────────────────
        print("\n   📱 生成竖版视频图片（9:16）...")
        images_dir_v = final_output_dir / "images_portrait"
        if images_dir_v.exists():
            shutil.rmtree(images_dir_v)
        images_dir_v.mkdir(exist_ok=True)
        for image_path in raw_all_images_v:
            shutil.copy2(image_path, images_dir_v / image_path.name)
        all_images_v = validate_image_count(images_dir_v, pages, "竖版")
        validate_portrait_images(images_dir_v, pages, "竖版")

        if logo_path:
            print("      → 遮盖竖版 NotebookLM Logo...")
            cover_v = cover_logo_on_images(images_dir_v, logo_path)
            if cover_v != len(all_images_v):
                print(f"   ⚠️ 竖版部分图片遮盖失败: {cover_v}/{len(all_images_v)}")
        else:
            print("      ⏭️ 跳过 Logo 遮盖（Logo 文件不存在）")
        print("      → 绘制竖版免责声明...")
        draw_disclaimer(images_dir_v, orientation="portrait")
        
        # ─── Step 9: 图片后处理完成 ─────────────────────────────────
        print("\n[9/10] 图片后处理已完成（PPTX 使用后处理前图片生成，不再用遮盖/免责声明后的图片重组）")
        
        # ─── Step 10: 清理 ──────────────────────────────────────────
        print("\n[10/10] 清理临时文件...")
        
        if not keep_temp and temp_dir and temp_dir.exists():
            shutil.rmtree(temp_dir)
            print(f"   ✅ 删除临时目录: {temp_dir}")
        
        # 删除笔记本
        if notebook_id:
            try:
                nb.delete_notebook(notebook_id)
                print(f"   ✅ 删除笔记本: {notebook_id[:8]}...")
            except Exception as e:
                print(f"   ⚠️ 删除笔记本失败: {e}")
        
        # ─── 页数校验 ───────────────────────────────────────────────────
        h_count = len(all_images_h)
        v_count = len(all_images_v)
        for w in warnings:
            print(f"   ⚠️ {w}")

        # ─── 完成 ───────────────────────────────────────────────────
        step_timer("done")
        total_sec = time.time() - run_start

        print("\n" + "="*60)
        print(f"✅ 流水线完成！（耗时 {total_sec:.0f}s，run_id={run_id}）")
        print(f"   📦 横版 PPTX: {final_pptx_h}")
        print(f"   📦 竖版 PPTX: {final_pptx_v}")
        print(f"   🖼️  横版图片: {images_dir_h} ({h_count}/{pages} 张)")
        print(f"   🖼️  竖版图片: {images_dir_v} ({v_count}/{pages} 张)")
        if warnings:
            print(f"   ⚠️ 警告: {len(warnings)} 项")
        print("="*60)

        result = {
            "status": "ok" if not warnings else "partial",
            "run_id": run_id,
            "total_seconds": round(total_sec, 1),
            "step_times": {name: round(sec, 1) for name, sec in step_times},
            "pptx_landscape": str(final_pptx_h),
            "pptx_portrait": str(final_pptx_v),
            "images_dir_landscape": str(images_dir_h),
            "images_dir_portrait": str(images_dir_v),
            "page_count_landscape": h_count,
            "page_count_portrait": v_count,
            "expected_pages": pages,
            "warnings": warnings,
            "notebook_id": notebook_id,
            "output_dir": str(final_output_dir),
        }

        # 保存运行结果 JSON
        result_file = final_output_dir / "run_result.json"
        result_file.write_text(json.dumps(result, ensure_ascii=False, indent=2))
        print(f"   📋 运行结果: {result_file}")

        return result
        
    except Exception as e:
        print(f"\n❌ 流水线失败: {e}")
        import traceback
        traceback.print_exc()
        
        # 清理
        if not keep_temp:
            if temp_dir and temp_dir.exists():
                shutil.rmtree(temp_dir)
            if notebook_id:
                try:
                    nb.delete_notebook(notebook_id)
                except:
                    pass
        
        sys.exit(1)

# ═══════════════════════════════════════════════════════════════════════
# CLI 入口
# ═══════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="NotebookLM MD → PPTX 完整流水线",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例:
    %(prog)s report.md
    %(prog)s report.md --title "我的报告" --pages 40
    %(prog)s report.md --output-dir ~/Desktop/output
    %(prog)s report.md --logo ~/path/to/logo.png
        """
    )
    
    parser.add_argument("md_file", help="Markdown 文件路径")
    parser.add_argument("--title", "-t", help="标题（默认使用文件名）")
    parser.add_argument("--pages", "-p", type=int, default=30, 
                        help="目标页数（默认 30）")
    parser.add_argument("--output-dir", "-o", type=Path,
                        help=f"输出目录（默认 {DEFAULT_OUTPUT_DIR}/<笔记本名称>）")
    parser.add_argument("--logo", "-l", type=Path,
                        help=f"Logo 底图路径（默认 {DEFAULT_LOGO_PATH}）")
    parser.add_argument("--keep-temp", action="store_true",
                        help="保留临时文件（调试用）")
    parser.add_argument("--skip-chrome-check", action="store_true",
                        help="跳过 Chrome 9222 远程调试检查（仅调试用）")
    parser.add_argument("--initial-interval", type=int, default=540,
                        help="初始等待时间/秒（默认 540）")
    parser.add_argument("--max-interval", type=int, default=60,
                        help="轮询间隔/秒（默认 60）")
    parser.add_argument("--timeout", type=int, default=1800,
                        help="总超时/秒（默认 1800）")
    
    args = parser.parse_args()
    
    # 验证文件
    md_path = Path(args.md_file).expanduser().resolve()
    if not md_path.exists():
        print(f"❌ 文件不存在: {md_path}")
        sys.exit(1)
    
    wait_config = {
        "initial_interval": args.initial_interval,
        "max_interval": args.max_interval,
        "timeout": args.timeout,
    }
    
    result = main(
        md_file=md_path,
        title=args.title,
        pages=args.pages,
        output_dir=args.output_dir,
        logo_path=args.logo,
        keep_temp=args.keep_temp,
        wait_config=wait_config,
        check_chrome=not args.skip_chrome_check,
    )
