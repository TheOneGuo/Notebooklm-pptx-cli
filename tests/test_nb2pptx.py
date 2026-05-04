from pathlib import Path
import importlib.util

from PIL import Image
import pytest


def load_module(module_name: str, relative_path: str):
    repo_root = Path(__file__).resolve().parents[1]
    module_path = repo_root / relative_path
    spec = importlib.util.spec_from_file_location(module_name, module_path)
    module = importlib.util.module_from_spec(spec)
    assert spec.loader is not None
    spec.loader.exec_module(module)
    return module


nb2pptx = load_module("nb2pptx", "scripts/nb2pptx.py")


def write_png(path: Path, size: tuple[int, int]):
    Image.new("RGB", size, "white").save(path)


def test_prepare_upload_source_copies_chinese_filename_to_ascii_temp(tmp_path):
    md_file = tmp_path / "中文研报.md"
    md_file.write_text("# 标题\n\n内容", encoding="utf-8")
    temp_dir = tmp_path / "temp"
    temp_dir.mkdir()

    upload_path = nb2pptx.prepare_upload_source(md_file, temp_dir)

    assert upload_path == temp_dir / "source.md"
    assert upload_path.read_text(encoding="utf-8") == md_file.read_text(encoding="utf-8")


def test_copy_images_with_expected_count_trims_extra_pages(tmp_path):
    src_dir = tmp_path / "src"
    src_dir.mkdir()
    extracted = []
    for idx in range(1, 7):
        path = src_dir / f"P{idx}.png"
        path.write_bytes(f"img-{idx}".encode("utf-8"))
        extracted.append(path)

    output_dir = tmp_path / "output"
    output_dir.mkdir()

    copied, warnings = nb2pptx.copy_images_with_expected_count(
        extracted,
        output_dir,
        start_page=4,
        expected_count=3,
        deck_label="竖版后半",
    )

    assert [path.name for path in copied] == ["P4.png", "P5.png", "P6.png"]
    assert [path.read_bytes() for path in copied] == [b"img-1", b"img-2", b"img-3"]
    assert warnings == ["竖版后半 图片数 6 > 预期 3，仅保留前 3 页"]


def test_copy_images_with_expected_count_raises_when_under_generated(tmp_path):
    src_dir = tmp_path / "src"
    src_dir.mkdir()
    extracted = []
    for idx in range(1, 3):
        path = src_dir / f"P{idx}.png"
        path.write_bytes(f"img-{idx}".encode("utf-8"))
        extracted.append(path)

    output_dir = tmp_path / "output"
    output_dir.mkdir()

    with pytest.raises(RuntimeError, match="竖版后半 图片数 2 < 预期 3"):
        nb2pptx.copy_images_with_expected_count(
            extracted,
            output_dir,
            start_page=4,
            expected_count=3,
            deck_label="竖版后半",
        )


def test_validate_portrait_images_detects_landscape_pages(tmp_path):
    images_dir = tmp_path / "images_portrait"
    images_dir.mkdir()
    write_png(images_dir / "P1.png", (768, 1376))
    write_png(images_dir / "P2.png", (768, 1376))
    write_png(images_dir / "P3.png", (1376, 768))

    with pytest.raises(RuntimeError, match="竖版页面方向校验失败") as excinfo:
        nb2pptx.validate_portrait_images(images_dir, expected_pages=3, deck_label="竖版前半")

    message = str(excinfo.value)
    assert "P3" in message
    assert "1376x768" in message


def test_create_pptx_from_images_preserves_image_aspect_ratio(tmp_path):
    from pptx import Presentation

    image_path = tmp_path / "P1.png"
    write_png(image_path, (1600, 900))
    output_path = tmp_path / "portrait.pptx"

    nb2pptx.create_pptx_from_images([image_path], output_path, orientation="portrait")

    prs = Presentation(str(output_path))
    slide = prs.slides[0]
    picture = slide.shapes[0]

    assert len(prs.slides) == 1
    assert picture.width / picture.height == pytest.approx(1600 / 900, rel=0.01)
    assert not (picture.width == prs.slide_width and picture.height == prs.slide_height)
    assert picture.height == prs.slide_height


def test_main_generates_merged_30_page_outputs_and_postprocessed_images(tmp_path, monkeypatch):
    md_file = tmp_path / "input.md"
    md_file.write_text("# 报告\n\n内容", encoding="utf-8")
    logo_path = tmp_path / "logo.png"
    write_png(logo_path, (120, 40))
    output_dir = tmp_path / "out"

    generated_ids = []
    generate_counter = {"value": 0}
    postprocess_calls = []

    class FakeNotebook:
        def __init__(self, cli_path):
            self.cli_path = cli_path

        def create_notebook(self, title):
            return "nb-1"

        def add_source_file(self, md_path, notebook_id):
            return "source-a", md_path.stem

        def rename_source(self, source_id, notebook_id, new_name):
            return None

        def run(self, args, timeout=30, check=True):
            if args[:1] == ["ask"] and "--save-as-note" not in args:
                return {"output": "测试报告"}
            if args[:1] == ["ask"] and "--save-as-note" in args:
                return {"output": "ok"}
            if args[:2] == ["notes", "list"]:
                return {"notes": [{"title": "4页PPT大纲", "id": "note-1"}]}
            if args[:2] == ["generate", "slide-deck"]:
                generate_counter["value"] += 1
                artifact_id = f"artifact-p{generate_counter['value']}"
                generated_ids.append(artifact_id)
                return {"task_id": artifact_id, "artifact_id": artifact_id}
            if args[:2] == ["artifact", "list"]:
                return {
                    "artifacts": [
                        {"id": artifact_id, "status": 3, "kind": "slide_deck"}
                        for artifact_id in generated_ids
                    ]
                }
            raise AssertionError(f"unexpected run args: {args}")

        def rename_notebook(self, notebook_id, title):
            return None

        def get_note_content(self, note_id, notebook_id):
            return "大纲内容" * 30

        def add_source_text(self, note_content, notebook_id, note_title):
            return "source-b"

        def download_slides(self, output_path, artifact_id, notebook_id, fmt="pptx"):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(artifact_id.encode("utf-8"))
            return output_path

        def delete_notebook(self, notebook_id):
            return None

    def fake_extract(pptx_path, extract_dir):
        extract_dir.mkdir(parents=True, exist_ok=True)
        paths = []
        for idx in range(1, 3):
            path = extract_dir / f"P{idx}.png"
            write_png(path, (768, 1376))
            paths.append(path)
        return paths

    monkeypatch.setattr(nb2pptx, "find_notebooklm_cli", lambda: Path("/fake/notebooklm"))
    monkeypatch.setattr(nb2pptx, "NotebookLM", FakeNotebook)
    monkeypatch.setattr(nb2pptx, "extract_images_from_pptx", fake_extract)
    monkeypatch.setattr(nb2pptx, "cover_logo_on_images", lambda *args, **kwargs: postprocess_calls.append("logo") or 4)
    monkeypatch.setattr(nb2pptx, "draw_disclaimer", lambda *args, **kwargs: postprocess_calls.append("disclaimer") or 4)

    result = nb2pptx.main(
        md_file=md_file,
        title="测试报告",
        pages=4,
        output_dir=output_dir,
        logo_path=logo_path,
        keep_temp=True,
        wait_config={"initial_interval": 0, "max_interval": 0, "timeout": 99},
        check_chrome=False,
    )

    assert result["page_count_landscape"] == 4
    assert result["page_count_portrait"] == 4
    assert nb2pptx.count_pptx_slides(output_dir / "测试报告_横版.pptx") == 4
    assert nb2pptx.count_pptx_slides(output_dir / "测试报告_竖版.pptx") == 4
    assert sorted(path.name for path in (output_dir / "images_landscape").glob("P*.png")) == ["P1.png", "P2.png", "P3.png", "P4.png"]
    assert sorted(path.name for path in (output_dir / "images_portrait").glob("P*.png")) == ["P1.png", "P2.png", "P3.png", "P4.png"]
    assert postprocess_calls == ["logo", "disclaimer", "logo", "disclaimer"]
